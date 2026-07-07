"""Tests for multi-format document loaders (PDF, DOCX, Excel)."""
import pytest


class TestLoadPdf:
    def test_load_pdf_extracts_text(self, tmp_path):
        """PDF with text content should be extracted."""
        from pypdf import PdfWriter
        from app.rag.loader import load_pdf

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        path = str(tmp_path / "test.pdf")
        with open(path, "wb") as f:
            writer.write(f)

        result = load_pdf(path)
        assert "content" in result
        assert "metadata" in result
        assert result["metadata"]["source"] == "test.pdf"
        assert result["metadata"]["file_type"] == "pdf"

    def test_load_docx_extracts_paragraphs(self, tmp_path):
        """DOCX paragraphs should be extracted as text."""
        from docx import Document
        from app.rag.loader import load_docx

        doc = Document()
        doc.add_paragraph("This is a test document with English content.")
        doc.add_paragraph("Second paragraph here.")
        path = str(tmp_path / "test.docx")
        doc.save(path)

        result = load_docx(path)
        assert "test document" in result["content"].lower()
        assert "second paragraph" in result["content"].lower()
        assert result["metadata"]["language"] == "en"
        assert result["metadata"]["file_type"] == "docx"

    def test_load_docx_chinese_content(self, tmp_path):
        """DOCX with Chinese content should detect language as zh."""
        from docx import Document
        from app.rag.loader import load_docx

        doc = Document()
        doc.add_paragraph("这是一份中文测试文档，用于验证语言检测功能。")
        path = str(tmp_path / "chinese.docx")
        doc.save(path)

        result = load_docx(path)
        assert result["metadata"]["language"] == "zh"

    def test_load_excel_converts_rows(self, tmp_path):
        """Excel rows should be converted to 'col: value' text."""
        from openpyxl import Workbook
        from app.rag.loader import load_excel

        wb = Workbook()
        ws = wb.active
        ws.append(["Name", "Age", "City"])
        ws.append(["Alice", 30, "Beijing"])
        ws.append(["Bob", 25, "Shanghai"])
        path = str(tmp_path / "test.xlsx")
        wb.save(path)

        result = load_excel(path)
        assert "Name" in result["content"]
        assert "Alice" in result["content"]
        assert "Beijing" in result["content"]
        assert result["metadata"]["file_type"] == "xlsx"

    def test_load_csv_converts_rows_with_metadata(self, tmp_path):
        from app.rag.loader import load_csv

        csv_path = tmp_path / "sales.csv"
        csv_path.write_text("region,revenue\nAPAC,1200\nEMEA,900", encoding="utf-8")

        result = load_csv(str(csv_path))

        assert "region: APAC" in result["content"]
        assert "revenue: 1200" in result["content"]
        assert result["metadata"]["file_type"] == "csv"
        assert result["metadata"]["source"] == "sales.csv"

    def test_csv_handles_bom_semicolon_and_blank_lines(self, tmp_path):
        from app.rag.ingestion.extractors import extract_csv_document

        csv_path = tmp_path / "pipeline.csv"
        csv_path.write_text("\ufeffstage;owner\nDiscovery;Sales\n\nDelivery;CS\n", encoding="utf-8")

        chunks = extract_csv_document(csv_path)

        assert len(chunks) == 1
        assert chunks[0].metadata["headers"] == ["stage", "owner"]
        assert chunks[0].metadata["delimiter"] == ";"
        assert chunks[0].metadata["row_start"] == 2
        assert chunks[0].metadata["row_end"] == 3
        assert "stage: Discovery" in chunks[0].text
        assert "owner: CS" in chunks[0].text

    def test_load_csv_without_header_raises_clear_error(self, tmp_path):
        from app.rag.loader import load_csv

        csv_path = tmp_path / "no-header.csv"
        csv_path.write_text("APAC,1200\nEMEA,900\n", encoding="utf-8")

        with pytest.raises(ValueError, match="CSV contains no header row"):
            load_csv(str(csv_path))

    def test_load_single_document_dispatches_docx(self, tmp_path):
        """load_single_document should handle .docx files."""
        from docx import Document
        from app.rag.loader import load_single_document

        doc = Document()
        doc.add_paragraph("Test content for dispatch.")
        path = str(tmp_path / "dispatch.docx")
        doc.save(path)

        result = load_single_document(path)
        assert result["metadata"]["source"] == "dispatch.docx"
        assert "dispatch" in result["content"].lower() or "test content" in result["content"].lower()
        assert result["metadata"]["uploaded"] is True

    def test_load_single_document_dispatches_xlsx(self, tmp_path):
        """load_single_document should handle .xlsx files."""
        from openpyxl import Workbook
        from app.rag.loader import load_single_document

        wb = Workbook()
        ws = wb.active
        ws.append(["Col1", "Col2"])
        ws.append(["A", "B"])
        path = str(tmp_path / "data.xlsx")
        wb.save(path)

        result = load_single_document(path)
        assert result["metadata"]["source"] == "data.xlsx"
        assert "A" in result["content"]

    def test_unsupported_extension_raises(self, tmp_path):
        """Unsupported file types should raise ValueError."""
        from app.rag.loader import load_single_document

        path = str(tmp_path / "test.exe")
        with open(path, "w") as f:
            f.write("a,b,c")

        with pytest.raises(ValueError, match="Unsupported"):
            load_single_document(path)

    def test_load_single_document_dispatches_pptx(self, tmp_path):
        from pptx import Presentation
        from app.rag.loader import load_single_document

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Quarterly Review"
        textbox = slide.shapes.add_textbox(0, 0, 3000000, 1000000)
        textbox.text_frame.text = "Revenue grew in APAC."
        path = tmp_path / "review.pptx"
        prs.save(path)

        result = load_single_document(str(path))

        assert result["metadata"]["source"] == "review.pptx"
        assert result["metadata"]["file_type"] == "pptx"
        assert "Quarterly Review" in result["content"]
        assert "APAC" in result["content"]

    def test_pptx_extracts_tables_and_notes(self, tmp_path):
        from pptx import Presentation
        from app.rag.ingestion.extractors import extract_pptx_document

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Pipeline Review"
        table = slide.shapes.add_table(2, 2, 0, 0, 3000000, 1000000).table
        table.cell(0, 0).text = "Stage"
        table.cell(0, 1).text = "Owner"
        table.cell(1, 0).text = "Discovery"
        table.cell(1, 1).text = "Sales"
        slide.notes_slide.notes_text_frame.text = "Mention rollout risk."
        path = tmp_path / "pipeline.pptx"
        prs.save(path)

        chunks = extract_pptx_document(path)

        assert chunks
        assert chunks[0].metadata["slide_number"] == 1
        assert chunks[0].metadata["slide_title"] == "Pipeline Review"
        assert "Stage | Owner" in chunks[0].text
        assert "Discovery | Sales" in chunks[0].text
        assert "Speaker notes" in chunks[0].text
        assert "rollout risk" in chunks[0].text

    def test_pptx_extracts_multiple_slides(self, tmp_path):
        from pptx import Presentation
        from app.rag.ingestion.extractors import extract_pptx_document

        prs = Presentation()
        first = prs.slides.add_slide(prs.slide_layouts[5])
        first.shapes.title.text = "Executive Summary"
        first.shapes.add_textbox(0, 0, 3000000, 1000000).text_frame.text = "North star metric improved."
        second = prs.slides.add_slide(prs.slide_layouts[5])
        second.shapes.title.text = "Risk Register"
        second.shapes.add_textbox(0, 0, 3000000, 1000000).text_frame.text = "Renewal risk requires owner."
        path = tmp_path / "multi-slide.pptx"
        prs.save(path)

        chunks = extract_pptx_document(path)

        assert [chunk.metadata["slide_number"] for chunk in chunks] == [1, 2]
        assert chunks[0].metadata["slide_title"] == "Executive Summary"
        assert chunks[1].metadata["slide_title"] == "Risk Register"
        assert "North star metric" in chunks[0].text
        assert "Renewal risk" in chunks[1].text

    def test_empty_pptx_raises_value_error(self, tmp_path):
        from pptx import Presentation
        from app.rag.loader import load_pptx

        prs = Presentation()
        path = tmp_path / "empty.pptx"
        prs.save(path)

        with pytest.raises(ValueError, match="no extractable text"):
            load_pptx(str(path))

    def test_docx_extracts_table_content(self, tmp_path):
        from docx import Document
        from app.rag.ingestion.extractors import extract_docx_document

        doc = Document()
        doc.add_heading("Pricing", level=1)
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Plan"
        table.cell(0, 1).text = "Price"
        table.cell(1, 0).text = "Pro"
        table.cell(1, 1).text = "$99"
        path = tmp_path / "pricing.docx"
        doc.save(path)

        chunks = extract_docx_document(path)

        assert any("Plan | Price" in chunk.text for chunk in chunks)
        assert any("Pro | $99" in chunk.text for chunk in chunks)
        assert any(chunk.metadata.get("element_type") == "table" for chunk in chunks)
        assert any(chunk.metadata.get("heading_path") == "Pricing" for chunk in chunks)

    def test_xlsx_preserves_sheet_and_row_metadata(self, tmp_path):
        from openpyxl import Workbook
        from app.rag.ingestion.extractors import extract_xlsx_document

        wb = Workbook()
        ws = wb.active
        ws.title = "Pipeline"
        ws.append(["Stage", "Owner"])
        ws.append(["Discovery", "Sales"])
        ws.append(["Delivery", "CS"])
        path = tmp_path / "pipeline.xlsx"
        wb.save(path)

        chunks = extract_xlsx_document(path)

        assert chunks
        assert chunks[0].metadata["sheet_name"] == "Pipeline"
        assert chunks[0].metadata["row_start"] == 2
        assert chunks[0].metadata["row_end"] == 3
        assert chunks[0].metadata["headers"] == ["Stage", "Owner"]
        assert "Stage: Discovery" in chunks[0].text

    def test_pdf_extracts_page_metadata_for_blank_pdf(self, tmp_path):
        from pypdf import PdfWriter
        from app.rag.ingestion.extractors import extract_pdf_document

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        path = tmp_path / "blank.pdf"
        with open(path, "wb") as f:
            writer.write(f)

        doc = extract_pdf_document(path)

        assert doc.metadata["file_type"] == "pdf"
        assert doc.metadata["page_count"] == 1
        assert "warnings" in doc.metadata
