"""Upload golden corpus and deterministic rank metrics.

This module is intentionally test-owned: it defines the fixed files and
sentinel queries used to prove that upload -> index -> search works end-to-end.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any


@dataclass(frozen=True)
class UploadGoldenCase:
    file_type: str
    filename: str
    sentinel: str
    expected_slug: str
    query: str


UPLOAD_GOLDEN_CASES = [
    UploadGoldenCase(
        file_type="csv",
        filename="aureon-golden-csv.csv",
        sentinel="AUREON_GOLDEN_SENTINEL_CSV_20260706",
        expected_slug="aureon-golden-csv",
        query="AUREON_GOLDEN_SENTINEL_CSV_20260706",
    ),
    UploadGoldenCase(
        file_type="docx",
        filename="aureon-golden-docx.docx",
        sentinel="AUREON_GOLDEN_SENTINEL_DOCX_20260706",
        expected_slug="aureon-golden-docx",
        query="AUREON_GOLDEN_SENTINEL_DOCX_20260706",
    ),
    UploadGoldenCase(
        file_type="md",
        filename="aureon-golden-md.md",
        sentinel="AUREON_GOLDEN_SENTINEL_MD_20260706",
        expected_slug="aureon-golden-md",
        query="AUREON_GOLDEN_SENTINEL_MD_20260706",
    ),
    UploadGoldenCase(
        file_type="pdf",
        filename="aureon-golden-pdf.pdf",
        sentinel="AUREON_GOLDEN_SENTINEL_PDF_20260706",
        expected_slug="aureon-golden-pdf",
        query="AUREON_GOLDEN_SENTINEL_PDF_20260706",
    ),
    UploadGoldenCase(
        file_type="pptx",
        filename="aureon-golden-pptx.pptx",
        sentinel="AUREON_GOLDEN_SENTINEL_PPTX_20260706",
        expected_slug="aureon-golden-pptx",
        query="AUREON_GOLDEN_SENTINEL_PPTX_20260706",
    ),
    UploadGoldenCase(
        file_type="txt",
        filename="aureon-golden-txt.txt",
        sentinel="AUREON_GOLDEN_SENTINEL_TXT_20260706",
        expected_slug="aureon-golden-txt",
        query="AUREON_GOLDEN_SENTINEL_TXT_20260706",
    ),
    UploadGoldenCase(
        file_type="xlsx",
        filename="aureon-golden-xlsx.xlsx",
        sentinel="AUREON_GOLDEN_SENTINEL_XLSX_20260706",
        expected_slug="aureon-golden-xlsx",
        query="AUREON_GOLDEN_SENTINEL_XLSX_20260706",
    ),
]


def create_upload_golden_files(output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    paths = {case.file_type: output_dir / case.filename for case in UPLOAD_GOLDEN_CASES}

    _write_csv(paths["csv"], _case("csv"))
    _write_docx(paths["docx"], _case("docx"))
    _write_md(paths["md"], _case("md"))
    _write_pdf(paths["pdf"], _case("pdf"))
    _write_pptx(paths["pptx"], _case("pptx"))
    _write_txt(paths["txt"], _case("txt"))
    _write_xlsx(paths["xlsx"], _case("xlsx"))
    return paths


def evaluate_upload_search_results(
    cases: list[UploadGoldenCase],
    results_by_type: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    details = []
    failures = []
    matched = 0
    rank1 = 0

    for case in cases:
        result = results_by_type.get(case.file_type, {})
        sources = result.get("sources", [])
        rank = _find_rank(case, sources)
        top = sources[0].get("slug") if sources else None

        if rank is not None:
            matched += 1
        if rank == 1:
            rank1 += 1
        if rank != 1:
            failures.append({
                "file_type": case.file_type,
                "sentinel": case.sentinel,
                "expected_slug": case.expected_slug,
                "rank": rank,
                "top": top,
            })

        details.append({
            "file_type": case.file_type,
            "sentinel": case.sentinel,
            "expected_slug": case.expected_slug,
            "rank": rank,
            "top": top,
            "source_count": len(sources),
        })

    total = len(cases)
    return {
        "total": total,
        "matched": matched,
        "rank1": rank1,
        "recall_at_k": matched / total if total else 0.0,
        "rank1_rate": rank1 / total if total else 0.0,
        "details": details,
        "failures": failures,
    }


def _case(file_type: str) -> UploadGoldenCase:
    return next(case for case in UPLOAD_GOLDEN_CASES if case.file_type == file_type)


def _find_rank(case: UploadGoldenCase, sources: list[dict[str, Any]]) -> int | None:
    for index, source in enumerate(sources, start=1):
        haystack = " ".join(
            str(source.get(field, ""))
            for field in ("title", "slug", "chunk", "chunk_text_snippet")
        )
        if case.expected_slug in haystack or case.sentinel in haystack:
            return index
    return None


def _body(case: UploadGoldenCase) -> str:
    return (
        f"Golden upload regression file for {case.file_type}.\n"
        f"Unique sentinel: {case.sentinel}.\n"
        "Customer proof: exact lookup must return this uploaded file at rank 1."
    )


def _write_csv(path: Path, case: UploadGoldenCase) -> None:
    path.write_text(
        "document_id,department,sentinel,decision\n"
        f"csv-001,finance,{case.sentinel},approve renewal\n",
        encoding="utf-8",
    )


def _write_docx(path: Path, case: UploadGoldenCase) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Aureon Upload Golden DOCX", level=1)
    doc.add_paragraph(_body(case))
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "field"
    table.cell(0, 1).text = "value"
    table.cell(1, 0).text = "sentinel"
    table.cell(1, 1).text = case.sentinel
    doc.save(str(path))


def _write_md(path: Path, case: UploadGoldenCase) -> None:
    path.write_text(
        f"---\ntitle: Aureon Golden MD\nslug: {case.expected_slug}\n---\n\n# Upload Golden MD\n\n{_body(case)}\n",
        encoding="utf-8",
    )


def _write_pdf(path: Path, case: UploadGoldenCase) -> None:
    text = _body(case).replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    chunks = [b"%PDF-1.4\n"]
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(sum(len(chunk) for chunk in chunks))
        chunks.append(f"{index} 0 obj\n".encode("ascii") + obj + b"\nendobj\n")
    xref_offset = sum(len(chunk) for chunk in chunks)
    chunks.append(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    chunks.append(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        chunks.append(f"{offset:010d} 00000 n \n".encode("ascii"))
    chunks.append(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode("ascii")
    )
    path.write_bytes(b"".join(chunks))


def _write_pptx(path: Path, case: UploadGoldenCase) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Aureon Upload Golden PPTX"
    slide.placeholders[1].text = _body(case)
    prs.save(str(path))


def _write_txt(path: Path, case: UploadGoldenCase) -> None:
    path.write_text(_body(case), encoding="utf-8")


def _write_xlsx(path: Path, case: UploadGoldenCase) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Golden"
    ws.append(["document_id", "sentinel", "owner"])
    ws.append(["xlsx-001", case.sentinel, "enterprise-success"])
    wb.save(str(path))
