"""File-type specific extraction helpers for ingestion."""

from __future__ import annotations

import csv
import re
from io import StringIO
from pathlib import Path

from app.rag.ingestion.models import ChunkRecord, IngestedDocument
from app.rag.ingestion.normalizer import normalize_text
from app.utils.lang_detect import detect_language as _detect_text_language


STRUCTURED_CHUNK_MAX_CHARS = 6000
STRUCTURED_CHUNK_MAX_ESTIMATED_TOKENS = 480


def _estimate_embedding_tokens(text: str) -> int:
    cjk_chars = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")
    other_chars = max(0, len(text) - cjk_chars)
    return cjk_chars + (max(1, other_chars // 4) if other_chars else 0)


def _pick_csv_delimiter(sample: str, fallback: str) -> str:
    first_line = next((line for line in sample.splitlines() if line.strip()), "")
    if first_line.count(";") > first_line.count(","):
        return ";"
    return fallback


def _looks_like_csv_header(sample: str, dialect: csv.Dialect, delimiter: str) -> bool:
    try:
        rows = list(csv.reader(StringIO(sample), dialect=dialect, delimiter=delimiter))
    except csv.Error:
        return False
    rows = [row for row in rows if any(cell.strip() for cell in row)]
    if len(rows) < 2:
        return False
    first = [cell.strip() for cell in rows[0]]
    if not first or any(not cell for cell in first):
        return False
    if len(set(first)) != len(first):
        return False
    if any(cell.replace(".", "", 1).isdigit() for cell in first):
        return False
    return True


def _split_long_line(
    line: str,
    max_chars: int = STRUCTURED_CHUNK_MAX_CHARS,
    max_estimated_tokens: int = STRUCTURED_CHUNK_MAX_ESTIMATED_TOKENS,
) -> list[str]:
    if len(line) <= max_chars and _estimate_embedding_tokens(line) <= max_estimated_tokens:
        return [line]
    parts = []
    start = 0
    while start < len(line):
        end = min(len(line), start + max_chars)
        while end > start + 1 and _estimate_embedding_tokens(line[start:end]) > max_estimated_tokens:
            span = end - start
            end = start + max(1, span * max_estimated_tokens // _estimate_embedding_tokens(line[start:end]))
        parts.append(line[start:end])
        start = end
    return parts


def parse_frontmatter(content: str) -> tuple[dict[str, object], str]:
    match = re.match(r"^---\s*\n(.*?)\n---\s*\n(.*)", content, re.DOTALL)
    if not match:
        return {}, content

    frontmatter_text = match.group(1)
    body = match.group(2).strip()

    metadata: dict[str, object] = {}
    for line in frontmatter_text.strip().split("\n"):
        if ":" in line:
            key, _, value = line.partition(":")
            key = key.strip()
            value = value.strip().strip('"').strip("'")
            if value.startswith("[") and value.endswith("]"):
                value = [v.strip().strip('"').strip("'") for v in value[1:-1].split(",")]
            metadata[key] = value
    return metadata, body


def extract_markdown_document(path: Path) -> IngestedDocument:
    content = path.read_text(encoding="utf-8")
    metadata, body = parse_frontmatter(content)
    lang = _detect_text_language(body[:500])
    return IngestedDocument(
        metadata={
            "source": path.name,
            "title": metadata.get("title", path.stem),
            "slug": metadata.get("slug", path.stem),
            "tags": metadata.get("tags", []),
            "category": metadata.get("category", ""),
            "filepath": str(path),
            "language": lang,
            "file_type": "md",
        },
        content=normalize_text(body),
    )


def extract_text_document(path: Path) -> IngestedDocument:
    content = normalize_text(path.read_text(encoding="utf-8"))
    return IngestedDocument(
        metadata={
            "source": path.name,
            "title": path.stem,
            "slug": path.stem,
            "tags": [],
            "category": "upload",
            "filepath": str(path),
            "language": _detect_text_language(content[:500]),
            "file_type": "txt",
        },
        content=content,
    )


def extract_pdf_document(path: Path) -> IngestedDocument:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text:
            pages.append(f"[Page {page_number}]\n{text.strip()}")
    content = normalize_text("\n\n".join(pages))
    warnings = []
    if len(content.strip()) < 20:
        warnings.append("PDF contains little or no extractable text; it may be scanned or image-based.")
    lang = _detect_text_language(content[:500]) if content else "en"
    return IngestedDocument(
        metadata={
            "source": path.name,
            "title": path.stem,
            "slug": path.stem,
            "tags": [],
            "category": "upload",
            "filepath": str(path),
            "language": lang,
            "file_type": "pdf",
            "page_count": len(reader.pages),
            "warnings": warnings,
        },
        content=content,
    )


def _make_upload_chunk(
    path: Path,
    content: str,
    file_type: str,
    element_type: str,
    heading_path: str = "",
) -> ChunkRecord:
    lang = _detect_text_language(content[:500]) if content else "en"
    metadata = {
        "source": path.name,
        "title": path.stem,
        "slug": path.stem,
        "tags": [],
        "category": "upload",
        "filepath": str(path),
        "language": lang,
        "file_type": file_type,
        "element_type": element_type,
    }
    if heading_path:
        metadata["heading_path"] = heading_path
    return ChunkRecord(text=content, metadata=metadata)


def _docx_table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    header = rows[0]
    lines = [" | ".join(header), " | ".join(["---"] * len(header))]
    for row in rows[1:]:
        lines.append(" | ".join(row))
    return "\n".join(lines)


def extract_docx_document(path: Path) -> list[ChunkRecord]:
    from docx import Document

    doc = Document(str(path))
    chunks: list[ChunkRecord] = []
    heading_stack: list[str] = []
    paragraph_parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        style_name = paragraph.style.name if paragraph.style else ""
        if style_name.startswith("Heading"):
            if paragraph_parts:
                content = normalize_text("\n\n".join(paragraph_parts))
                chunks.append(_make_upload_chunk(path, content, "docx", "paragraph", " > ".join(heading_stack)))
                paragraph_parts = []

            raw_level = style_name.replace("Heading", "").strip()
            level = int(raw_level or "1") if raw_level.isdigit() else 1
            heading_stack = heading_stack[: max(level - 1, 0)]
            heading_stack.append(text)

        paragraph_parts.append(text)

    if paragraph_parts:
        content = normalize_text("\n\n".join(paragraph_parts))
        chunks.append(_make_upload_chunk(path, content, "docx", "paragraph", " > ".join(heading_stack)))

    for index, table in enumerate(doc.tables, start=1):
        table_text = normalize_text(_docx_table_to_markdown(table))
        if table_text:
            chunk = _make_upload_chunk(path, table_text, "docx", "table", " > ".join(heading_stack))
            chunk.metadata["table_index"] = index
            chunks.append(chunk)

    return chunks


def extract_csv_document(path: Path) -> list[ChunkRecord]:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        content = f.read()

    sample = content.lstrip("\ufeff")[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample) if sample else csv.excel
    except csv.Error:
        dialect = csv.excel
    delimiter = _pick_csv_delimiter(sample, dialect.delimiter)

    try:
        has_header = csv.Sniffer().has_header(sample) if sample else False
    except csv.Error:
        has_header = False
    has_header = has_header or _looks_like_csv_header(sample, dialect, delimiter)
    if not has_header:
        return []

    reader = csv.DictReader(StringIO(content), dialect=dialect, delimiter=delimiter)
    headers = [str(header) for header in (reader.fieldnames or []) if header is not None]
    chunks: list[ChunkRecord] = []
    lines = []
    row_count = 0
    row_start = 2
    row_end = 1
    current_chars = 0
    current_tokens = 0
    header_line = _format_csv_header(headers, delimiter)
    header_tokens = _estimate_embedding_tokens(header_line)
    max_row_chars = max(1, STRUCTURED_CHUNK_MAX_CHARS - len(header_line) - 1)
    max_row_tokens = max(1, STRUCTURED_CHUNK_MAX_ESTIMATED_TOKENS - header_tokens)

    for row in reader:
        row_count += 1
        file_row_number = row_count + 1
        values = []
        for header in headers:
            value = row.get(header)
            if value not in (None, ""):
                values.append(str(value).strip().replace("\r", " ").replace("\n", " "))
            else:
                values.append("")
        if any(values):
            for line_part in _split_long_line(_format_csv_row(values, delimiter), max_row_chars, max_row_tokens):
                line_len = len(line_part) + 1
                line_tokens = _estimate_embedding_tokens(line_part)
                would_exceed_chars = len(header_line) + 1 + current_chars + line_len > STRUCTURED_CHUNK_MAX_CHARS
                would_exceed_tokens = header_tokens + current_tokens + line_tokens > STRUCTURED_CHUNK_MAX_ESTIMATED_TOKENS
                if lines and (would_exceed_chars or would_exceed_tokens):
                    chunks.append(_make_csv_chunk(path, headers, row_start, row_end, lines, delimiter))
                    lines = []
                    current_chars = 0
                    current_tokens = 0
                if not lines:
                    row_start = file_row_number
                lines.append(line_part)
                current_chars += line_len
                current_tokens += line_tokens
                row_end = file_row_number

    if lines:
        row_end = row_count + 1
        chunks.append(_make_csv_chunk(path, headers, row_start, row_end, lines, delimiter))

    return chunks


def _format_csv_header(headers: list[str], delimiter: str) -> str:
    separator = f"{delimiter} " if delimiter != "\t" else "\t"
    return "Columns: " + separator.join(headers)


def _format_csv_row(values: list[str], delimiter: str) -> str:
    separator = f"{delimiter} " if delimiter != "\t" else "\t"
    return separator.join(values)


def _make_csv_chunk(
    path: Path,
    headers: list[str],
    row_start: int,
    row_end: int,
    lines: list[str],
    delimiter: str,
) -> ChunkRecord:
    text = normalize_text("\n".join([_format_csv_header(headers, delimiter), *lines]))
    lang = _detect_text_language(text[:500]) if text else "en"
    return ChunkRecord(
        text=text,
        metadata={
            "source": path.name,
            "title": path.stem,
            "slug": path.stem,
            "tags": [],
            "category": "upload",
            "filepath": str(path),
            "language": lang,
            "file_type": "csv",
            "headers": headers,
            "row_start": row_start,
            "row_end": row_end,
            "delimiter": delimiter,
        },
    )


def _pptx_table_to_text(table) -> str:
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_pptx_document(path: Path) -> list[ChunkRecord]:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[ChunkRecord] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        parts = []
        slide_title = ""
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            slide_title = title_shape.text.strip().splitlines()[0] if title_shape.text.strip() else ""

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text = shape.text.strip()
                if not slide_title:
                    slide_title = text.splitlines()[0]
                parts.append(text)

            if getattr(shape, "has_table", False):
                table_text = _pptx_table_to_text(shape.table)
                if table_text:
                    parts.append(table_text)

        if getattr(slide, "has_notes_slide", False):
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"Speaker notes:\n{notes}")

        content = normalize_text("\n\n".join(parts))
        if not content:
            continue

        lang = _detect_text_language(content[:500])
        chunks.append(
            ChunkRecord(
                text=content,
                metadata={
                    "source": path.name,
                    "title": path.stem,
                    "slug": path.stem,
                    "tags": [],
                    "category": "upload",
                    "filepath": str(path),
                    "language": lang,
                    "file_type": "pptx",
                    "slide_number": slide_index,
                    "slide_title": slide_title,
                    "element_type": "slide",
                },
            )
        )

    return chunks


def extract_xlsx_document(path: Path) -> list[ChunkRecord]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[ChunkRecord] = []
    rows_per_chunk = 50

    for ws in wb.worksheets:
        headers: list[str] = []
        lines: list[str] = []
        row_start = 2
        row_end = 1
        current_chars = 0
        current_tokens = 0

        for row_number, row in enumerate(ws.iter_rows(values_only=True), start=1):
            if row_number == 1:
                headers = [str(h).strip() if h is not None else f"col{i}" for i, h in enumerate(row)]
                continue

            parts = []
            for header, value in zip(headers, row):
                if value is not None and str(value).strip():
                    parts.append(f"{header}: {value}")
            if parts:
                for line_part in _split_long_line(", ".join(parts)):
                    line_len = len(line_part) + 1
                    line_tokens = _estimate_embedding_tokens(line_part)
                    would_exceed_chars = current_chars + line_len > STRUCTURED_CHUNK_MAX_CHARS
                    would_exceed_tokens = current_tokens + line_tokens > STRUCTURED_CHUNK_MAX_ESTIMATED_TOKENS
                    if lines and (len(lines) >= rows_per_chunk or would_exceed_chars or would_exceed_tokens):
                        chunks.append(_make_xlsx_chunk(path, ws.title, headers, row_start, row_end, lines))
                        lines = []
                        current_chars = 0
                        current_tokens = 0
                    if not lines:
                        row_start = row_number
                    lines.append(line_part)
                    current_chars += line_len
                    current_tokens += line_tokens
                    row_end = row_number

            if len(lines) >= rows_per_chunk:
                chunks.append(_make_xlsx_chunk(path, ws.title, headers, row_start, row_end, lines))
                lines = []
                current_chars = 0
                current_tokens = 0

        if lines:
            chunks.append(_make_xlsx_chunk(path, ws.title, headers, row_start, row_end, lines))

    wb.close()
    return chunks


def _make_xlsx_chunk(
    path: Path,
    sheet_name: str,
    headers: list[str],
    row_start: int,
    row_end: int,
    lines: list[str],
) -> ChunkRecord:
    content = normalize_text("\n".join(lines))
    lang = _detect_text_language(content[:500]) if content else "en"
    return ChunkRecord(
        text=content,
        metadata={
            "source": path.name,
            "title": path.stem,
            "slug": path.stem,
            "tags": [],
            "category": "upload",
            "filepath": str(path),
            "language": lang,
            "file_type": "xlsx",
            "element_type": "table_rows",
            "sheet_name": sheet_name,
            "headers": headers,
            "row_start": row_start,
            "row_end": row_end,
        },
    )
